from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import os

class Element:
    def __init__(self, content, style, bounding_box):
        self.content = content
        self.style = style
        self.bounding_box = bounding_box
    
    def apply_font_style_on_run(self, run):
        if 'font_name' in self.style:
            run.font.name = self.style['font_name']
        if 'font_size' in self.style:
            run.font.size = Pt(self.style['font_size'])
        if 'font_color' in self.style:
            run.font.color.rgb = RGBColor(self.style['font_color']['r'], self.style['font_color']['g'], self.style['font_color']['b'])
        if 'bold' in self.style:
            run.font.bold = self.style['bold']
        if 'italics' in self.style:
            run.font.italic = self.style['italics']
        if 'underlined' in self.style:
            run.font.underline = self.style['underlined']

    
    def apply_font_style(self, shape):
        if 'font_name' in self.style:
            shape.text_frame.paragraphs[0].font.name = self.style['font_name']
        if 'font_size' in self.style:
            shape.text_frame.paragraphs[0].font.size = Pt(self.style['font_size'])
        if 'font_color' in self.style:
            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(self.style['font_color']['r'], self.style['font_color']['g'], self.style['font_color']['b'])
        if 'bold' in self.style:
            shape.text_frame.paragraphs[0].font.bold = self.style['bold']
        else:
            shape.text_frame.paragraphs[0].font.bold = False 
        if 'italics' in self.style:
            shape.text_frame.paragraphs[0].font.italic = self.style['italics']
        else:
            shape.text_frame.paragraphs[0].font.italic = False 
        if 'underlined' in self.style:
            shape.text_frame.paragraphs[0].font.underline = self.style['underlined']
        else:
            shape.text_frame.paragraphs[0].font.underline = False
    
    def position_element(self, shape):
        shape.left = Inches(self.bounding_box[0])
        shape.top = Inches(self.bounding_box[1])
        shape.width = Inches(self.bounding_box[2])
        shape.height = Inches(self.bounding_box[3])

        if shape.has_text_frame:
            shape.text_frame.margin_left = Pt(10)
            shape.text_frame.margin_right = Pt(10)
            shape.text_frame.margin_top = Pt(10)
            shape.text_frame.margin_bottom = Pt(10)


    def render(self, slide):
        raise NotImplementedError("Subclasses must implement this method")

class Title(Element):
    def __init__(self, content, style, bounding_box):
        super().__init__(content, style, bounding_box)
    def render(self, slide):
        title_shape = slide.shapes.title
        title_shape.text = self.content
        self.apply_font_style(title_shape)
        self.position_element(title_shape)
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def clean_up(self, slide):
        pass

class Description(Element):
    def __init__(self, content, style, bounding_box):
        super().__init__(content, style, bounding_box)

    def render(self, slide):
        left, top, width, height = self.bounding_box
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        textbox.text = self.content
        self.apply_font_style(textbox)
        textbox.text_frame.auto_size = True
        textbox.text_frame.word_wrap = True 
        self.position_element(textbox)
        self.textbox = textbox

class Enumeration(Description):
    def __init__(self, content, style, bounding_box, heading):
        super().__init__(content, style, bounding_box)
        self.heading = heading

    def render(self, slide):
        left, top, width, height = self.heading['xmin'], self.heading['ymin'], self.heading['width'], self.heading['height']
        enum_heading = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        enum_heading.text = self.heading['value']
        self.apply_font_style(enum_heading)
        enum_heading.text_frame.auto_size = True
        enum_heading.text_frame.word_wrap = True
        enum_shape = slide.shapes.placeholders[1]
        enum_tf = enum_shape.text_frame
        if self.content != []:            
            enum_tf.text = self.content[0]
            self.apply_font_style(enum_shape)
            for i, pt_text in enumerate(self.content):
                if i>0:
                    if isinstance(pt_text, str):
                        p = enum_tf.add_paragraph()
                        run = p.add_run()
                        run.text = pt_text
                        self.apply_font_style_on_run(run)

                    elif isinstance(pt_text, list):
                        for sub_pt in pt_text:
                            s_p = enum_tf.add_paragraph()
                            s_p.level = 1
                            sub_run = s_p.add_run()
                            sub_run.text = sub_pt
                            self.apply_font_style_on_run(run)
                    else:
                        raise Exception("Invalid Enumeration format")
                
        enum_tf.auto_size = True
        enum_tf.word_wrap = True

        self.position_element(enum_shape)
        self.enum_tf = enum_tf
            

class Figure(Element):
    def __init__(self, content, style, bounding_box, caption):
        super().__init__(content, style, bounding_box)
        self.caption = caption
        
    def render(self, slide):
        left, top, width, height = self.bounding_box
        img = slide.shapes.add_picture(self.content, Inches(left), Inches(top), Inches(width), Inches(height))
        left, top, width, height = self.caption['xmin'], self.caption['ymin'], self.caption['width'], self.caption['height']

        cap_shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        cap_shape.text = self.caption['value']
        self.apply_font_style(cap_shape)
        cap_shape.text_frame.auto_size = True
        cap_shape.text_frame.word_wrap = True
        cap_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # original_width, original_height = img.image.size
        # # Convert dimensions from pixels to inches
        # img.width = Inches(original_width * scale_factor / img.image.dpi[0])
        # img.height = Inches(original_height * scale_factor / img.image.dpi[1])
        self.image = img 

class Equation(Element):
    def __init__(self, content, style, bounding_box):
        super().__init__(content, style, bounding_box)

    def render(self, slide):
        left, top, width, height = self.bounding_box
        img = slide.shapes.add_picture(self.content, Inches(left), Inches(top), Inches(width), Inches(height))
        # original_width, original_height = img.image.size
        # # Convert dimensions from pixels to inches
        # img.width = Inches(original_width * scale_factor / img.image.dpi[0])
        # img.height = Inches(original_height * scale_factor / img.image.dpi[1])
        self.image = img 

class Table(Element):
    def __init__(self, content, style, bounding_box):
        super().__init__(content, style, bounding_box)

    def render(self, slide):
        left, top, width, height = self.bounding_box
        img = slide.shapes.add_picture(self.content, Inches(left), Inches(top), Inches(width), Inches(height))
        # original_width, original_height = img.image.size
        # # Convert dimensions from pixels to inches
        # img.width = Inches(original_width * scale_factor / img.image.dpi[0])
        # img.height = Inches(original_height * scale_factor / img.image.dpi[1])
        self.image = img 


class Footer(Element):
    def __init__(self, content, style, bounding_box, location):
        super().__init__(content, style, bounding_box)
        self.location = location
    def render(self, slide):
        left, top, width, height = self.bounding_box
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        textbox.text = self.content
        self.apply_font_style(textbox)
        textbox.text_frame.auto_size = True
        textbox.text_frame.word_wrap = True
        if self.location == 1:
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        elif self.location == 3:
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        else:
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        self.position_element(textbox)
        self.textbox = textbox


class PresentationGenerator:
    def __init__(self, json_payload, presentation_id):
        self.json_payload = json_payload
        self.presentation_id = presentation_id
        if os.path.exists(os.path.join('ppts', f'{presentation_id}.pptx')):
            self.presentation = Presentation(os.path.join('ppts', f'{presentation_id}.pptx'))
        else:
            self.presentation = Presentation()
            self.insert_title_slide()
    
    def insert_title_slide(self):
        title_font = self.json_payload["slides"][0]["elements"]["title"][0]["style"]["font_name"]
        title_slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        title_shape = title_slide.shapes.title
        title_shape.text = self.json_payload['topic']
        title_shape.left = Inches(0.5)
        title_shape.top = Inches(2)
        title_shape.width = Inches(9)
        title_shape.height = Inches(1.25)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.size = Pt(48)
        title_shape.text_frame.paragraphs[0].font.name = title_font
        presenter = title_slide.shapes.add_textbox(Inches(3.5),Inches(3.25),Inches(3),Inches(0.75))
        presenter.text = self.json_payload["presenter"]
        presenter.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        presenter.text_frame.paragraphs[0].font.italic = True
        presenter.text_frame.paragraphs[0].font.size = Pt(28)
        presenter.text_frame.paragraphs[0].font.name = title_font
        date = title_slide.shapes.add_textbox(Inches(4),Inches(4),Inches(2),Inches(0.5))
        date.text = self.json_payload["date"]
        date.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        date.text_frame.paragraphs[0].font.name = title_font

    def generate_presentation(self):
        slide_info = self.json_payload['slides'][-1]
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        if slide_info['bg_color']:
            slide.background.fill.fore_color.rgb = RGBColor(slide_info['bg_color']['r'], slide_info['bg_color']['g'], slide_info['bg_color']['b'])
        else:
            slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

        for element_type, elements in slide_info['elements'].items():
            for element_info in elements:
                if element_type == 'figures':
                    element = Figure(element_info['path'], element_info['caption']['style'], (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']), element_info['caption'])
                elif element_type == 'equations':
                    element = Equation(element_info['path'], None, (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']))
                elif element_type == 'tables':
                    element = Table(element_info['path'], None, (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']))
                elif element_type == 'description':
                    if element_info['label'] == "enumeration":
                        element = Enumeration(element_info['value'], element_info['style'], (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']), element_info['heading'])
                    else:
                        element = Description(element_info['value'], element_info['style'], (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']))
                elif element_type == 'title':
                    element = Title(element_info['value'], element_info['style'], (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']))
                elif element_type == 'footer':
                    element = Footer(element_info['value'], element_info['style'], (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']), element_info['location'])
                else:
                    raise ValueError(f"Unsupported element type: {element_type}")

                element.render(slide)


    
        # Clean-up empty elements
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                # print(shape.shape_type)
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    if not shape.text_frame.text:
                        sp = slide.shapes._spTree
                        sp.remove(shape._element)
        
        ppts_path = "ppts"
        if os.path.isdir(ppts_path) == False:
            os.mkdir(ppts_path)

        self.presentation.save(os.path.join(ppts_path, f'{self.presentation_id}.pptx'))
        os.startfile(os.path.join(ppts_path, f'{self.presentation_id}.pptx'))

def load_json_payload(file_path):
    with open(file_path, 'r') as file:
        return json.load(file)

def main():
    # Load the JSON payload
    print("\nGenerating final presentation...")
    buffer_folder_path = "output"
    base_topic_folder_path = "json"

    json_files = [f for f in os.listdir(buffer_folder_path) if f.endswith('.json')]

    for json_file in json_files:
        #Load JSON file from buffer
        presentation_id , _ = os.path.splitext(json_file)
        json_file_path = os.path.join(buffer_folder_path, json_file)
        json_payload = load_json_payload(json_file_path)
        # presentation_id = json_payload["n_slides"]
        #Generate Presentation from JSON file and save it
        presentation_generator = PresentationGenerator(json_payload, presentation_id)
        try:
            presentation_generator.generate_presentation()
            print(f"🟢 (1/1) Generated New presentation: {presentation_id}.")
        except:
            print(f"🔴 ERROR: in generating presentation")
        
    
        #Move JSON file from buffer to respective topic folder
        if not os.path.exists(base_topic_folder_path):
            os.makedirs(base_topic_folder_path)

        # destination_file_path = os.path.join(base_topic_folder_path, json_file)
        # os.rename(json_file_path, destination_file_path)
        # print(f"Moved {json_file} to {base_topic_folder_path}.")
        # print('\n')
    
if __name__ == "__main__":
    main()